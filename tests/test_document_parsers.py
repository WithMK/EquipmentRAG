from __future__ import annotations

import tempfile
import unittest
from base64 import b64decode
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from app.config import DocumentConfig
from app.parsers.document_parser import DocumentParseError
from app.parsers.document_parsers import DocumentParserRegistry
from app.parsers.document_scanner import DocumentScanner


class DocumentParserTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.config = DocumentConfig(
            enabled=True,
            source_paths=(self.root,),
            extensions=(".docx", ".pdf", ".pptx", ".xlsx", ".md", ".txt"),
            exclude_directories=("archive",),
            chunk_size=500,
            chunk_overlap=50,
            collection_name="document_chunks",
        )

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def test_scanner_hashes_files_and_loads_revision_sidecar(self) -> None:
        path = self.root / "Loader.md"
        path.write_text("# Loader\n\nVacuum interlock", encoding="utf-8")
        path.with_name("Loader.md.metadata.yaml").write_text(
            "document_id: loader-spec\n"
            "revision: Rev.3\n"
            "document_status: active\n"
            "is_latest: true\n",
            encoding="utf-8",
        )
        archive = self.root / "archive"
        archive.mkdir()
        (archive / "old.txt").write_text("obsolete", encoding="utf-8")

        sources = DocumentScanner(self.config).scan()

        self.assertEqual(len(sources), 1)
        self.assertEqual(sources[0].relative_path, "Loader.md")
        self.assertEqual(len(sources[0].file_hash), 64)
        self.assertEqual(sources[0].metadata["revision"], "Rev.3")

    def test_markdown_preserves_headings_lists_and_code_blocks(self) -> None:
        path = self.root / "Manual.md"
        path.write_text(
            "# Trimming\n\n## Loader\n\n- Vacuum sensor ON\n"
            "- Door closed\n\n```csharp\nVacuum.On();\n```\n",
            encoding="utf-8",
        )
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        self.assertEqual(document.title, "Trimming")
        self.assertEqual(
            [block.type for block in document.blocks],
            ["heading", "heading", "list", "code"],
        )
        self.assertIn("Vacuum.On", document.blocks[-1].text)

    def test_txt_uses_section_patterns_before_paragraph_fallback(self) -> None:
        path = self.root / "Trouble.txt"
        path.write_text(
            "1. Loader Alarm\nVacuum sensor was not detected.\n\n"
            "ACTION:\nCheck the vacuum sensor.",
            encoding="utf-8",
        )
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        self.assertEqual(document.blocks[0].type, "heading")
        self.assertEqual(document.blocks[1].type, "paragraph")
        self.assertEqual(document.blocks[2].text, "ACTION")

    def test_docx_preserves_heading_paragraph_list_and_table(self) -> None:
        from docx import Document

        path = self.root / "Specification.docx"
        document = Document()
        document.core_properties.title = "Loader Specification"
        document.add_heading("Loader Interlock", level=1)
        document.add_paragraph("Vacuum sensor must be ON.")
        document.add_paragraph("Door must be closed.", style="List Bullet")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Signal"
        table.cell(0, 1).text = "Condition"
        table.cell(1, 0).text = "Vacuum"
        table.cell(1, 1).text = "ON"
        document.save(path)
        source = DocumentScanner(self.config).scan()[0]

        normalized = DocumentParserRegistry().parse(source)

        self.assertEqual(normalized.title, "Loader Specification")
        self.assertEqual(
            [block.type for block in normalized.blocks],
            ["heading", "paragraph", "list", "table"],
        )
        self.assertEqual(normalized.blocks[-1].rows[1], ("Vacuum", "ON"))

    def test_text_pdf_preserves_page_and_removes_repeated_margins(self) -> None:
        path = self.root / "Alarm.pdf"
        _write_pdf(
            path,
            [
                ["EQUIPMENT MANUAL", "1 Loader", "Vacuum timeout", "CONFIDENTIAL"],
                ["EQUIPMENT MANUAL", "2 Recovery", "Check sensor", "CONFIDENTIAL"],
            ],
        )
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        text = "\n".join(block.text for block in document.blocks)
        self.assertNotIn("EQUIPMENT MANUAL", text)
        self.assertNotIn("CONFIDENTIAL", text)
        self.assertEqual({block.page for block in document.blocks}, {1, 2})
        self.assertIn("Vacuum timeout", text)

    def test_scanned_pdf_uses_configured_ocr_for_sparse_pages(self) -> None:
        path = self.root / "Scanned.pdf"
        _write_pdf(path, [[""]])
        source = DocumentScanner(self.config).scan()[0]
        ocr = _FakeOcr("스캔 문서 진공 센서 점검 절차")
        pixmap = SimpleNamespace(tobytes=lambda _format: b"rendered-png")
        page = SimpleNamespace(get_pixmap=lambda **_kwargs: pixmap)
        document = SimpleNamespace(load_page=lambda _index: page, close=lambda: None)
        fake_fitz = SimpleNamespace(
            open=lambda _path: document,
            Matrix=lambda x, y: (x, y),
        )

        with patch.dict("sys.modules", {"fitz": fake_fitz}):
            normalized = DocumentParserRegistry(
                ocr=ocr,
                pdf_ocr=True,
            ).parse(source)

        self.assertIn("진공 센서", "\n".join(block.text for block in normalized.blocks))
        self.assertEqual(ocr.calls, [(b"rendered-png", ".png")])

    def test_pptx_preserves_slides_tables_and_speaker_notes(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "DesignReview.pptx"
        presentation = Presentation()
        presentation.core_properties.title = "Loader Design Review"
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Loader Safety"
        body = slide.placeholders[1].text_frame
        body.text = "Safety door must be closed."
        bullet = body.add_paragraph()
        bullet.text = "Vacuum sensor must be ON."
        bullet.level = 1
        table = slide.shapes.add_table(
            2,
            2,
            Inches(1),
            Inches(4),
            Inches(6),
            Inches(1),
        ).table
        table.cell(0, 0).text = "Signal"
        table.cell(0, 1).text = "Condition"
        table.cell(1, 0).text = "Vacuum"
        table.cell(1, 1).text = "ON"
        slide.notes_slide.notes_text_frame.text = "Verify PLC input X100."
        presentation.save(path)
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        self.assertEqual(document.title, "Loader Design Review")
        self.assertTrue(all(block.slide == 1 for block in document.blocks))
        self.assertEqual(document.blocks[0].text, "Loader Safety")
        self.assertTrue(any(block.type == "table" for block in document.blocks))
        self.assertTrue(
            any("Verify PLC input X100" in block.text for block in document.blocks)
        )

    def test_pptx_image_ocr_adds_traceable_text(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        image_path = self.root / "alarm.png"
        image_path.write_bytes(
            b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
                "/x8AAusB9Wl2nKsAAAAASUVORK5CYII="
            )
        )
        path = self.root / "ImageReview.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
        presentation.save(path)
        source = DocumentScanner(self.config).scan()[0]
        ocr = _FakeOcr("AL101 진공 타임아웃")

        document = DocumentParserRegistry(
            ocr=ocr,
            pptx_image_ocr=True,
        ).parse(source)

        self.assertTrue(any(block.text == "Image OCR" for block in document.blocks))
        self.assertTrue(any("AL101" in block.text for block in document.blocks))
        self.assertTrue(all(block.slide == 1 for block in document.blocks))

    def test_xlsx_preserves_sheets_regions_formulas_and_cell_ranges(self) -> None:
        from openpyxl import Workbook

        path = self.root / "Signals.xlsx"
        workbook = Workbook()
        workbook.properties.title = "Loader Signal List"
        sheet = workbook.active
        sheet.title = "Loader IO"
        sheet.append(("Signal", "Description"))
        sheet.append(("X100", "Vacuum sensor"))
        sheet.append(("Result", "=COUNTA(A2:A2)"))
        sheet.append((None, None))
        sheet.append(("Alarm", "Recovery"))
        sheet.append(("AL101", "Check vacuum"))
        workbook.save(path)
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        self.assertEqual(document.title, "Loader Signal List")
        tables = [block for block in document.blocks if block.type == "table"]
        self.assertEqual([block.cell_range for block in tables], ["A1:B3", "A5:B6"])
        self.assertTrue(all(block.sheet == "Loader IO" for block in tables))
        self.assertIn("Formula: =COUNTA(A2:A2)", tables[0].text)
        self.assertIn("AL101", tables[1].text)

    def test_xlsx_extracts_chart_titles_series_and_references(self) -> None:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference

        path = self.root / "AlarmChart.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Alarms"
        sheet.append(("Alarm", "Count"))
        sheet.append(("AL101", 5))
        sheet.append(("AL102", 2))
        chart = BarChart()
        chart.title = "Alarm Count"
        chart.x_axis.title = "Alarm Code"
        chart.y_axis.title = "Occurrences"
        chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True)
        chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=3))
        sheet.add_chart(chart, "D2")
        workbook.save(path)
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry(
            xlsx_chart_extraction=True,
        ).parse(source)

        text = "\n".join(block.text for block in document.blocks)
        self.assertIn("Chart: Alarm Count", text)
        self.assertIn("X axis: Alarm Code", text)
        self.assertIn("Y axis: Occurrences", text)
        self.assertIn("'Alarms'!$B$2:$B$3", text)
        self.assertIn("'Alarms'!$A$2:$A$3", text)

    def test_pptx_reads_grouped_text_and_titleless_slides(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        path = self.root / "GroupedReview.pptx"
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[1])
        first.shapes.title.text = "Sequence Review"
        first.placeholders[1].text = "Loader pickup sequence"
        second = presentation.slides.add_slide(presentation.slide_layouts[6])
        group = second.shapes.add_group_shape()
        text_box = group.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
        )
        text_box.text = "Manual recovery requires safety reset."
        presentation.save(path)
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        headings = [block for block in document.blocks if block.type == "heading"]
        self.assertEqual([(block.text, block.slide) for block in headings], [
            ("Sequence Review", 1),
            ("Slide 2", 2),
        ])
        self.assertTrue(
            any(
                block.slide == 2 and "Manual recovery" in block.text
                for block in document.blocks
            )
        )

    def test_xlsx_reads_multiple_sheets_and_typed_values(self) -> None:
        from datetime import date

        from openpyxl import Workbook

        path = self.root / "Parameters.xlsx"
        workbook = Workbook()
        parameters = workbook.active
        parameters.title = "Parameters"
        parameters.append(("Enabled", "Limit", "Effective Date"))
        parameters.append((True, 12.5, date(2026, 8, 24)))
        alarms = workbook.create_sheet("Alarm List")
        alarms.append(("Code", "Description"))
        alarms.append(("AL101", "Vacuum timeout"))
        workbook.save(path)
        source = DocumentScanner(self.config).scan()[0]

        document = DocumentParserRegistry().parse(source)

        tables = [block for block in document.blocks if block.type == "table"]
        self.assertEqual(
            [(block.sheet, block.cell_range) for block in tables],
            [("Parameters", "A1:C2"), ("Alarm List", "A1:B2")],
        )
        self.assertIn("TRUE", tables[0].text)
        self.assertIn("12.5", tables[0].text)
        self.assertIn("2026-08-24", tables[0].text)

    def test_rejects_corrupt_office_files_with_clear_errors(self) -> None:
        pptx = self.root / "Broken.pptx"
        pptx.write_bytes(b"not a presentation")
        source = DocumentScanner(self.config).scan()[0]

        with self.assertRaisesRegex(DocumentParseError, "Unable to parse PPTX"):
            DocumentParserRegistry().parse(source)

        pptx.unlink()
        xlsx = self.root / "Broken.xlsx"
        xlsx.write_bytes(b"not a workbook")
        source = DocumentScanner(self.config).scan()[0]

        with self.assertRaisesRegex(DocumentParseError, "Unable to parse XLSX"):
            DocumentParserRegistry().parse(source)

    def test_rejects_xlsx_above_the_sheet_scan_limit(self) -> None:
        from openpyxl import Workbook

        path = self.root / "Oversized.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.cell(row=2000, column=501).value = "outside safety limit"
        workbook.save(path)
        source = DocumentScanner(self.config).scan()[0]

        with self.assertRaisesRegex(DocumentParseError, "too large to scan safely"):
            DocumentParserRegistry().parse(source)


def _write_pdf(path: Path, page_lines: list[list[str]]) -> None:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    for lines in page_lines:
        page = writer.add_blank_page(width=612, height=792)
        resources = DictionaryObject(
            {
                NameObject("/Font"): DictionaryObject(
                    {NameObject("/F1"): font_ref}
                )
            }
        )
        page[NameObject("/Resources")] = resources
        commands = ["BT", "/F1 12 Tf", "72 740 Td", "14 TL"]
        for index, line in enumerate(lines):
            escaped = line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
            if index:
                commands.append("T*")
            commands.append(f"({escaped}) Tj")
        commands.append("ET")
        stream = DecodedStreamObject()
        stream.set_data("\n".join(commands).encode("ascii"))
        page[NameObject("/Contents")] = writer._add_object(stream)
    with path.open("wb") as output:
        writer.write(output)


class _FakeOcr:
    def __init__(self, result: str) -> None:
        self.result = result
        self.calls: list[tuple[bytes, str]] = []

    def recognize(self, image_bytes: bytes, *, extension: str = ".png") -> str:
        self.calls.append((image_bytes, extension))
        return self.result


if __name__ == "__main__":
    unittest.main()
