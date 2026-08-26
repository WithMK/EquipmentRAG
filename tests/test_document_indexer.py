from __future__ import annotations

import tempfile
import unittest
from dataclasses import replace
from pathlib import Path
from typing import Any, Sequence

from app.config import (
    AppConfig,
    ChromaConfig,
    DocumentConfig,
    EmbeddingConfig,
    EquipmentConfig,
    LlmConfig,
    LoggingConfig,
    SearchConfig,
    SourceConfig,
)
from app.document_indexer import DocumentIndexerError, IncrementalDocumentIndexer
from app.retrieval.sqlite_fts import SQLiteFtsStore
from app.vectorstore.chroma_store import VectorRecord


class FakeEmbedding:
    dimension = 3

    def __init__(self) -> None:
        self.calls: list[list[str]] = []
        self.fail = False

    def embed_documents(self, texts: Sequence[str]) -> list[list[float]]:
        values = list(texts)
        self.calls.append(values)
        if self.fail:
            raise RuntimeError("synthetic embedding failure")
        return [[1.0, float(index % 2), 0.0] for index, _ in enumerate(values)]


class FakeStore:
    def __init__(self) -> None:
        self.records: dict[str, VectorRecord] = {}
        self.delete_filters: list[dict[str, Any]] = []
        self.delete_id_calls: list[tuple[str, ...]] = []

    def upsert(self, records: Sequence[VectorRecord]) -> int:
        values = list(records)
        for record in values:
            self.records[record.id] = record
        return len(values)

    def delete_by_ids(self, ids: Sequence[str]) -> int:
        values = tuple(ids)
        self.delete_id_calls.append(values)
        deleted = 0
        for record_id in values:
            if self.records.pop(record_id, None) is not None:
                deleted += 1
        return deleted

    def delete_where(self, where: dict[str, Any]) -> int:
        self.delete_filters.append(where)
        deleted = [
            record_id
            for record_id, record in self.records.items()
            if all(getattr(record.metadata, key) == value for key, value in where.items())
        ]
        for record_id in deleted:
            del self.records[record_id]
        return len(deleted)


class IncrementalDocumentIndexerTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.documents = self.root / "documents"
        self.documents.mkdir()
        self.embedding = FakeEmbedding()
        self.store = FakeStore()
        self.state_path = self.root / "chroma" / "document-state.json"
        self.config = self._config()

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _config(self) -> AppConfig:
        return AppConfig(
            project_root=self.root,
            equipment=EquipmentConfig(name="Trimming"),
            source=SourceConfig(
                path=self.root / "source",
                include_extensions=(".cs",),
                exclude_directories=("bin", "obj"),
                chunk_size=1000,
                chunk_overlap=100,
            ),
            embedding=EmbeddingConfig(
                model_path=self.root / "model",
                batch_size=8,
                device="cpu",
                normalize_embeddings=True,
            ),
            chromadb=ChromaConfig(
                path=self.root / "chroma",
                collection_name="equipment_code",
            ),
            search=SearchConfig(top_k=3),
            llm=LlmConfig(
                provider="llama_cpp",
                base_url="http://127.0.0.1:8080/v1",
                model="local-model",
                request_timeout_seconds=30,
            ),
            logging=LoggingConfig(level="INFO", path=self.root / "rag.log"),
            document=DocumentConfig(
                enabled=True,
                source_paths=(self.documents,),
                extensions=(".md", ".txt"),
                exclude_directories=("archive",),
                chunk_size=300,
                chunk_overlap=30,
                collection_name="document_chunks",
            ),
        )

    def _indexer(
        self, config: AppConfig | None = None
    ) -> IncrementalDocumentIndexer:
        return IncrementalDocumentIndexer(
            config or self.config,
            embedding=self.embedding,
            vector_store=self.store,
            state_path=self.state_path,
        )

    def test_indexes_incrementally_with_metadata_and_deletions(self) -> None:
        spec = self.documents / "Loader.md"
        spec.write_text("# Loader\n\nVacuum interlock condition", encoding="utf-8")
        sidecar = spec.with_name("Loader.md.metadata.yaml")
        sidecar.write_text(
            "document_id: loader-spec\n"
            "unit: Loader\n"
            "document_type: Specification\n"
            "revision: Rev.3\n",
            encoding="utf-8",
        )
        trouble = self.documents / "Trouble.txt"
        trouble.write_text("ALARM:\nVacuum timeout", encoding="utf-8")

        first = self._indexer().run()

        self.assertEqual(first.new_files, ("Loader.md", "Trouble.txt"))
        self.assertEqual(self.store.delete_filters, [{"equipment": "Trimming"}])
        self.assertGreater(first.upserted_chunks, 0)
        metadata = next(
            record.metadata
            for record in self.store.records.values()
            if record.metadata.document_id == "loader-spec"
        )
        self.assertEqual(metadata.unit, "Loader")
        self.assertEqual(metadata.revision, "Rev.3")
        self.assertEqual(metadata.source_type, "document")
        embedding_calls = len(self.embedding.calls)

        unchanged = self._indexer().run()

        self.assertEqual(unchanged.skipped_files, ("Loader.md", "Trouble.txt"))
        self.assertEqual(len(self.embedding.calls), embedding_calls)

        sidecar.write_text(
            "document_id: loader-spec\n"
            "unit: Loader\n"
            "document_type: Specification\n"
            "revision: Rev.4\n",
            encoding="utf-8",
        )
        trouble.unlink()
        (self.documents / "Review.md").write_text(
            "# Design Review\n\nVacuum input review action",
            encoding="utf-8",
        )
        incremental = self._indexer().run()

        self.assertEqual(incremental.new_files, ("Review.md",))
        self.assertEqual(incremental.changed_files, ("Loader.md",))
        self.assertEqual(incremental.deleted_files, ("Trouble.txt",))
        self.assertEqual(
            {
                record.metadata.revision
                for record in self.store.records.values()
                if record.metadata.document_id == "loader-spec"
            },
            {"Rev.4"},
        )
        self.assertTrue(
            any(
                record.metadata.file_name == "Review.md"
                for record in self.store.records.values()
            )
        )
        self.assertNotIn("Vacuum timeout", [r.document for r in self.store.records.values()])

    def test_dry_run_does_not_embed_or_write_state(self) -> None:
        (self.documents / "Manual.md").write_text(
            "# Manual\n\nLoader operation",
            encoding="utf-8",
        )

        report = self._indexer().run(dry_run=True)

        self.assertTrue(report.dry_run)
        self.assertGreater(report.prepared_chunks, 0)
        self.assertEqual(self.embedding.calls, [])
        self.assertFalse(self.state_path.exists())
        self.assertEqual(self.store.records, {})

    def test_embedding_failure_preserves_store_and_state(self) -> None:
        path = self.documents / "Manual.md"
        path.write_text("# Manual\n\nLoader operation", encoding="utf-8")
        self._indexer().run()
        state_before = self.state_path.read_bytes()
        records_before = dict(self.store.records)
        path.write_text("# Manual\n\nChanged operation", encoding="utf-8")
        self.embedding.fail = True

        with self.assertRaisesRegex(DocumentIndexerError, "provider failed"):
            self._indexer().run()

        self.assertEqual(self.state_path.read_bytes(), state_before)
        self.assertEqual(self.store.records, records_before)

    def test_full_and_settings_change_reindex_all_current_documents(self) -> None:
        (self.documents / "Manual.md").write_text(
            "# Manual\n\nLoader operation",
            encoding="utf-8",
        )
        (self.documents / "Alarm.txt").write_text(
            "ALARM:\nVacuum timeout",
            encoding="utf-8",
        )
        self._indexer().run()

        full = self._indexer().run(full_reindex=True)

        self.assertTrue(full.full_reindex)
        self.assertEqual(full.reindex_reason, "requested_full_reindex")
        self.assertEqual(full.reindexed_files, ("Alarm.txt", "Manual.md"))

        assert self.config.document is not None
        changed_config = replace(
            self.config,
            document=replace(self.config.document, chunk_size=250),
        )
        settings = self._indexer(changed_config).run(dry_run=True)

        self.assertTrue(settings.full_reindex)
        self.assertEqual(settings.reindex_reason, "index_settings_changed")
        self.assertEqual(settings.reindexed_files, ("Alarm.txt", "Manual.md"))

    def test_synchronizes_document_revision_to_sqlite_fts(self) -> None:
        manual = self.documents / "Manual.md"
        manual.write_text("# Alarm\n\nVacuum recovery procedure", encoding="utf-8")
        sidecar = manual.with_name("Manual.md.metadata.yaml")
        sidecar.write_text("revision: Rev.1\n", encoding="utf-8")
        lexical = SQLiteFtsStore(self.root / "lexical.sqlite3")
        config = replace(
            self.config,
            search=replace(
                self.config.search,
                lexical_backend="sqlite_fts5",
                lexical_path=lexical.path,
            ),
        )
        indexer = IncrementalDocumentIndexer(
            config,
            embedding=self.embedding,
            vector_store=self.store,
            lexical_store=lexical,
            state_path=self.state_path,
        )

        indexer.run()
        self.assertEqual(
            len(lexical.search("Vacuum", 5, filters={"revision": "Rev.1"})),
            1,
        )

        sidecar.write_text("revision: Rev.2\n", encoding="utf-8")
        indexer.run()
        self.assertEqual(
            lexical.search("Vacuum", 5, filters={"revision": "Rev.1"}),
            [],
        )
        self.assertEqual(
            len(lexical.search("Vacuum", 5, filters={"revision": "Rev.2"})),
            1,
        )
        lexical.close()

    def test_indexes_office_documents_with_exact_source_locations(self) -> None:
        from openpyxl import Workbook
        from pptx import Presentation

        presentation_path = self.documents / "Review.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Loader Safety"
        slide.placeholders[1].text = "Door closed and vacuum sensor ON"
        presentation.save(presentation_path)

        workbook_path = self.documents / "Signals.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Loader IO"
        sheet.append(("Signal", "Description"))
        sheet.append(("X100", "Vacuum sensor"))
        workbook.save(workbook_path)

        assert self.config.document is not None
        office_config = replace(
            self.config,
            document=replace(
                self.config.document,
                extensions=(".pptx", ".xlsx"),
            ),
        )

        report = self._indexer(office_config).run()

        self.assertEqual(report.new_files, ("Review.pptx", "Signals.xlsx"))
        metadata = [record.metadata for record in self.store.records.values()]
        self.assertTrue(
            any(item.file_name == "Review.pptx" and item.slide == 1 for item in metadata)
        )
        self.assertTrue(
            any(
                item.file_name == "Signals.xlsx"
                and item.sheet == "Loader IO"
                and item.cell_range == "A1:B2"
                for item in metadata
            )
        )


if __name__ == "__main__":
    unittest.main()
