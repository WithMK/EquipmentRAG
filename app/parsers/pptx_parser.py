"""PPTX parser with slide, table, and speaker-note traceability."""

from __future__ import annotations

from datetime import date, datetime
from typing import Any, Iterable

from app.models.document_models import (
    DocumentBlock,
    DocumentSourceFile,
    NormalizedDocument,
)
from app.parsers.document_parser import (
    DocumentParseError,
    build_normalized_document,
)


class PptxDocumentParser:
    def parse(self, source: DocumentSourceFile) -> NormalizedDocument:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise DocumentParseError(
                "python-pptx is required to parse PPTX documents"
            ) from exc

        try:
            presentation = Presentation(str(source.path))
            blocks, first_title = self._blocks(presentation, MSO_SHAPE_TYPE)
        except Exception as exc:
            raise DocumentParseError(
                f"Unable to parse PPTX: {source.path}"
            ) from exc

        properties = presentation.core_properties
        return build_normalized_document(
            source,
            blocks,
            detected_title=_string(properties.title) or first_title,
            detected_created_date=_date_string(properties.created),
        )

    def _blocks(self, presentation: Any, shape_types: Any) -> tuple[list[DocumentBlock], str]:
        blocks: list[DocumentBlock] = []
        first_title = ""
        for slide_number, slide in enumerate(presentation.slides, 1):
            title_shape = slide.shapes.title
            title = _shape_text(title_shape) or f"Slide {slide_number}"
            if not first_title and title_shape is not None:
                first_title = title
            blocks.append(
                DocumentBlock(
                    "heading",
                    title,
                    level=1,
                    slide=slide_number,
                )
            )

            title_id = getattr(title_shape, "shape_id", None)
            for shape in _iter_shapes(slide.shapes, shape_types):
                if getattr(shape, "shape_id", None) == title_id:
                    continue
                if getattr(shape, "has_table", False):
                    rows = _table_rows(shape.table)
                    if rows:
                        blocks.append(
                            DocumentBlock(
                                "table",
                                _format_rows(rows),
                                slide=slide_number,
                                rows=rows,
                            )
                        )
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    block_type = "list" if paragraph.level > 0 else "paragraph"
                    blocks.append(
                        DocumentBlock(
                            block_type,
                            text,
                            slide=slide_number,
                        )
                    )

            notes = _speaker_notes(slide)
            if notes:
                blocks.append(
                    DocumentBlock(
                        "heading",
                        "Speaker Notes",
                        level=2,
                        slide=slide_number,
                    )
                )
                blocks.append(
                    DocumentBlock(
                        "paragraph",
                        notes,
                        slide=slide_number,
                    )
                )
        return blocks, first_title


def _iter_shapes(shapes: Iterable[Any], shape_types: Any) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == shape_types.GROUP:
            yield from _iter_shapes(shape.shapes, shape_types)
        else:
            yield shape


def _shape_text(shape: Any) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def _table_rows(table: Any) -> tuple[tuple[str, ...], ...]:
    rows = tuple(
        tuple(cell.text.strip() for cell in row.cells)
        for row in table.rows
    )
    return tuple(row for row in rows if any(row))


def _format_rows(rows: tuple[tuple[str, ...], ...]) -> str:
    return "\n".join(" | ".join(row) for row in rows)


def _speaker_notes(slide: Any) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        frame = slide.notes_slide.notes_text_frame
    except (AttributeError, ValueError):
        return ""
    return "\n".join(
        paragraph.text.strip()
        for paragraph in frame.paragraphs
        if paragraph.text.strip()
    )


def _string(value: Any) -> str:
    return value.strip() if isinstance(value, str) else ""


def _date_string(value: Any) -> str:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return ""
